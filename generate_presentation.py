import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # 16:9 widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette: Deep Space Navy, Tech Blue, Emerald Green, Amber, Red, White, Slate
    BG_COLOR = RGBColor(10, 14, 26)           # #0a0e1a
    CARD_BG = RGBColor(17, 24, 43)            # #11182b
    CARD_BORDER = RGBColor(35, 49, 80)        # #233150
    ACCENT_BLUE = RGBColor(75, 107, 245)      # #4b6bf5
    ACCENT_CYAN = RGBColor(56, 189, 248)      # #38bdf8
    ACCENT_GREEN = RGBColor(16, 185, 129)     # #10b981
    ACCENT_AMBER = RGBColor(245, 158, 11)     # #f59e0b
    ACCENT_RED = RGBColor(239, 68, 68)        # #ef4444
    TEXT_WHITE = RGBColor(248, 250, 252)      # #f8fafc
    TEXT_MUTED = RGBColor(148, 163, 184)      # #94a3b8

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_header(slide, title_text, category_text="SMART INDIA HACKATHON 2026"):
        # Category Tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(11.7), Inches(0.35))
        tf_c = cat_box.text_frame
        tf_c.word_wrap = True
        p_c = tf_c.paragraphs[0]
        p_c.text = category_text.upper()
        p_c.font.size = Pt(11)
        p_c.font.bold = True
        p_c.font.color.rgb = ACCENT_CYAN

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11.7), Inches(0.8))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(25)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_WHITE

    blank_layout = prs.slide_layouts[6]

    # ==========================================
    # SLIDE 0: TITLE SLIDE
    # ==========================================
    slide0 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide0)

    center_card = slide0.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.0), Inches(0.9), Inches(11.333), Inches(5.7)
    )
    center_card.fill.solid()
    center_card.fill.fore_color.rgb = CARD_BG
    center_card.line.color.rgb = ACCENT_BLUE
    center_card.line.width = Pt(2)

    badge_box = slide0.shapes.add_textbox(Inches(1.5), Inches(1.2), Inches(10.3), Inches(0.4))
    tf_b0 = badge_box.text_frame
    p_b0 = tf_b0.paragraphs[0]
    p_b0.text = "SMART INDIA HACKATHON 2026 • AI & SATELLITE ATMOSPHERIC INTELLIGENCE"
    p_b0.font.size = Pt(12)
    p_b0.font.bold = True
    p_b0.font.color.rgb = ACCENT_GREEN

    t_box = slide0.shapes.add_textbox(Inches(1.5), Inches(1.65), Inches(10.3), Inches(1.5))
    tf_t0 = t_box.text_frame
    p_t0 = tf_t0.paragraphs[0]
    p_t0.text = "VayuShetra: Atmospheric Intelligence Platform"
    p_t0.font.size = Pt(36)
    p_t0.font.bold = True
    p_t0.font.color.rgb = TEXT_WHITE

    sub_p = tf_t0.add_paragraph()
    sub_p.text = "Hyperlocal Air Quality Forecasting, Sentinel-5P HCHO Hotspot Tracking & Source Attribution"
    sub_p.font.size = Pt(17)
    sub_p.font.color.rgb = ACCENT_CYAN

    # 4 Highlights
    highlights = [
        ("🛰️ HCHO & FIRMS Hotspots", "Detects invisible smoldering crop residue fires via Sentinel-5P HCHO"),
        ("🌪️ Lagrangian Wind AI", "Physics-driven smoke plume trajectory simulation across 300+ km corridors"),
        ("⚡ 48-72h Early Warning", "Pre-emptive district intelligence replacing retroactive city shutdowns"),
        ("🎯 Chemical Source Attribution", "Decomposes air pollution into Farm Fire vs Vehicle vs Factory percentages")
    ]
    for i, (h_title, h_desc) in enumerate(highlights):
        left = Inches(1.5 + i * 2.58)
        h_card = slide0.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(3.85), Inches(2.42), Inches(2.3))
        h_card.fill.solid()
        h_card.fill.fore_color.rgb = RGBColor(13, 19, 34)
        h_card.line.color.rgb = CARD_BORDER
        h_card.line.width = Pt(1)

        htf = h_card.text_frame
        htf.word_wrap = True
        hp1 = htf.paragraphs[0]
        hp1.text = h_title
        hp1.font.size = Pt(12)
        hp1.font.bold = True
        hp1.font.color.rgb = TEXT_WHITE

        hp2 = htf.add_paragraph()
        hp2.text = "\n" + h_desc
        hp2.font.size = Pt(10)
        hp2.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 1: PROBLEM STATEMENT
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)
    add_header(slide1, "1. Problem Statement: Scale of Crisis & Critical Blindspots")

    p_cards = [
        ("Health & Economic Catastrophe",
         "• 1.67 Million premature deaths in India annually linked to air pollution (Lancet).\n• ₹2.7 Lakh Crore ($36.8B) annual economic loss from worker morbidity & healthcare load.\n• Post-monsoon AQI reaches 450+ (25x to 30x above WHO limits), shutting schools & industries.",
         ACCENT_RED),
        ("Massive Ground Sensor Gap",
         "• Ground stations (CPCB) cost ₹1.5–2 Crore each to deploy + ₹30 Lakh/year maintenance.\n• 90% of stations are clustered in metro cities; 600,000+ rural villages have ZERO sensors.\n• Agricultural stubble burning belts in Punjab/Haryana/UP operate in total data blindness.",
         ACCENT_AMBER),
        ("Invisible Fires & Blame Games",
         "• Smoldering crop fires emit high Formaldehyde (HCHO) but low heat, evading standard thermal cameras.\n• Authorities lack causal source attribution (Farm Fires vs. Traffic vs. Industry).\n• Results in retroactive, blanket city shutdowns instead of targeted, pre-emptive enforcement.",
         ACCENT_BLUE)
    ]

    for i, (ctitle, cbody, ccolor) in enumerate(p_cards):
        left = Inches(0.8 + i * 3.95)
        card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.6), Inches(3.75), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = ccolor
        card.line.width = Pt(1.5)

        ctf = card.text_frame
        ctf.word_wrap = True
        cp1 = ctf.paragraphs[0]
        cp1.text = ctitle
        cp1.font.size = Pt(15)
        cp1.font.bold = True
        cp1.font.color.rgb = ccolor

        cp2 = ctf.add_paragraph()
        cp2.text = "\n" + cbody
        cp2.font.size = Pt(11.5)
        cp2.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 2: PROPOSED SOLUTION
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "2. Proposed Solution: VayuShetra Atmospheric Intelligence")

    sol_cards = [
        ("🛰️ Space-Borne Sensing with Sentinel-5P HCHO",
         "Leverages Sentinel-5P TROPOMI Formaldehyde (HCHO) columnar density and NASA VIIRS 375m hotspots to track active & smoldering biomass burning across every square kilometer of India.",
         ACCENT_AMBER),
        ("🌪️ Lagrangian Wind Transport Modeling",
         "Integrates ERA5 meteorological wind kinematics with Lagrangian plume dispersion to model exact smoke trajectories traveling 300+ km from northern farms to urban basins.",
         ACCENT_GREEN),
        ("⚡ 48-72h AI Predictive Forecasting",
         "Multi-seasonal LightGBM and XGBoost machine learning architectures predicting district-level PM2.5, PM10, and AQI up to 3 days in advance with weather-adjusted feature weighting.",
         ACCENT_BLUE),
        ("🎯 Trace Gas Chemical Source Apportionment",
         "Computes columnar gas ratios (NO2/CO, HCHO/SO2) to scientifically separate regional pollution into exact percentage shares: Biomass Burning, Vehicular Exhaust, and Industrial Stacks.",
         ACCENT_CYAN)
    ]

    for i, (stitle, sbody, scolor) in enumerate(sol_cards):
        row = i // 2
        col = i % 2
        left = Inches(0.8 + col * 5.95)
        top = Inches(1.6 + row * 2.65)
        
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.75), Inches(2.4))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = scolor
        card.line.width = Pt(1.5)

        stf = card.text_frame
        stf.word_wrap = True
        sp1 = stf.paragraphs[0]
        sp1.text = stitle
        sp1.font.size = Pt(15)
        sp1.font.bold = True
        sp1.font.color.rgb = scolor

        sp2 = stf.add_paragraph()
        sp2.text = "\n" + sbody
        sp2.font.size = Pt(11.5)
        sp2.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 3: KEY FEATURES & WORKFLOW
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "3. Key Features & Operational Workflow")

    workflow_steps = [
        ("Step 1: Ingest & Detect",
         "• Ingests Sentinel-5P trace gases (HCHO, NO2, SO2, CO) + NASA FIRMS thermal hotspots.\n• Pulls live CPCB ground station telemetry + ERA5 wind vector fields.\n• Cleans, spatial-aligns, and interpolates data across 15+ benchmark districts.",
         ACCENT_BLUE),
        ("Step 2: Model & Apportion",
         "• DBSCAN Spatial Clustering groups isolated fire pixels into regional burning centers.\n• Lagrangian Plume Dispersion tracks downwind smoke transport.\n• Chemical Fingerprinting computes real-time source apportionment percentages.",
         ACCENT_CYAN),
        ("Step 3: Forecast & Act",
         "• 48-72h LightGBM/XGBoost generates district-level AQI forecasts.\n• Automated GRAP Advisory Engine maps forecast to Stage I–IV mandates.\n• Delivers actionable alerts to district collectors, health officials, and citizens.",
         ACCENT_GREEN)
    ]

    for i, (wtitle, wbody, wcolor) in enumerate(workflow_steps):
        left = Inches(0.8 + i * 3.95)
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.6), Inches(3.75), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = wcolor
        card.line.width = Pt(1.5)

        wtf = card.text_frame
        wtf.word_wrap = True
        wp1 = wtf.paragraphs[0]
        wp1.text = wtitle
        wp1.font.size = Pt(15)
        wp1.font.bold = True
        wp1.font.color.rgb = wcolor

        wp2 = wtf.add_paragraph()
        wp2.text = "\n" + wbody
        wp2.font.size = Pt(11.5)
        wp2.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 4: TECHNOLOGY STACK
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "4. Technology Stack: Cloud-Native & Production-Ready")

    stack_cards = [
        ("🛰️ Space & Satellite Data Pipeline",
         "• Sentinel-5P TROPOMI (NetCDF4 Products: HCHO, NO2, SO2, CO, AOD)\n• NASA FIRMS 375m VIIRS & MODIS Thermal Anomaly Streams\n• ERA5 & Open-Meteo Planetary Boundary Layer (PBL) & Wind Vectors (u, v)\n• CPCB National CAAQMS Ground Sensor Real-Time API Feeds",
         ACCENT_CYAN),
        ("🧠 AI, Machine Learning & Analytics",
         "• LightGBM & XGBoost (Multi-Seasonal 48h-72h Predictive Models)\n• Scikit-Learn (DBSCAN Spatial Clustering for Hotspot Identification)\n• SciPy & NumPy (Lagrangian Plume Dispersion & Partial Correlation)\n• Pandas & NetCDF4 Engine (Vectorized Geospatial Data Processing)",
         ACCENT_BLUE),
        ("⚡ Backend & Data Infrastructure",
         "• Python 3.10+ & FastAPI (High-performance Asynchronous Microservices)\n• SQLite / PostgreSQL Database with spatial caching & session indexing\n• Automated CRON schedulers for continuous 15-minute satellite sync\n• Robust Offline Fallbacks with spatial k-NN interpolation",
         ACCENT_GREEN),
        ("🖥️ Frontend & UI Experience",
         "• React 18 + Vite (Ultra-fast responsive single-page application)\n• TailwindCSS & Custom Glassmorphism Theme (Dark mode, neon accents)\n• Recharts & Canvas 2D (Interactive spatial charts & vector wind flows)\n• Lucide Icons & Responsive District Analytics Spatial Matrix",
         ACCENT_AMBER)
    ]

    for i, (stitle, sbody, scolor) in enumerate(stack_cards):
        row = i // 2
        col = i % 2
        left = Inches(0.8 + col * 5.95)
        top = Inches(1.6 + row * 2.65)
        
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.75), Inches(2.4))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = scolor
        card.line.width = Pt(1.5)

        stf = card.text_frame
        stf.word_wrap = True
        sp1 = stf.paragraphs[0]
        sp1.text = stitle
        sp1.font.size = Pt(14)
        sp1.font.bold = True
        sp1.font.color.rgb = scolor

        sp2 = stf.add_paragraph()
        sp2.text = "\n" + sbody
        sp2.font.size = Pt(10.5)
        sp2.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 5: INNOVATION & IMPACT
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, "5. Innovation & Impact: Transforming Environmental Governance")

    impact_cards = [
        ("Key Technical Innovations",
         "• Sentinel-5P HCHO Chemical Tracer: Uncovers invisible smoldering fires that evade thermal infrared satellites.\n• Lagrangian Wind Transport AI: Maps interstate smoke travel (Punjab/Haryana -> Delhi) with kinematic accuracy.\n• Causal Source Apportionment: Mathematically decomposes pollution into Biomass vs Vehicle vs Factory shares.\n• Zero Ground Sensor Capex: 100% spatial coverage of India without spending ₹1.5 Cr per sensor station.",
         ACCENT_CYAN),
        ("Measurable Socio-Economic Impact",
         "• 48-72h Pre-emptive Action: Enables district collectors to deploy bio-decomposers & balers BEFORE burning occurs.\n• Health Protection for 50M+ Citizens: Delivers pollutant-specific health alerts for schools, hospitals, and high-risk patients.\n• Ending Interstate Blame Games: Provides unforgeable, audit-grade satellite evidence for policy enforcement.\n• Accelerated NCAP Progress: Directly aids National Clean Air Programme goal of 40% PM2.5 reduction by 2026.",
         ACCENT_GREEN)
    ]

    for i, (ititle, ibody, icolor) in enumerate(impact_cards):
        left = Inches(0.8 + i * 5.95)
        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.6), Inches(5.75), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = icolor
        card.line.width = Pt(1.5)

        itf = card.text_frame
        itf.word_wrap = True
        p1 = itf.paragraphs[0]
        p1.text = ititle
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = icolor

        p2 = itf.add_paragraph()
        p2.text = "\n" + ibody
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 6: MARKET POTENTIAL & SUSTAINABILITY
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_header(slide6, "6. Market Potential, Scalability & Sustainability")

    market_points = [
        ("🎯 Target Customers & Stakeholders",
         "• Government Bodies: Central Pollution Control Board (CPCB), State PCBs (Punjab, Haryana, Delhi, UP), MoEFCC.\n• District Administrations: 700+ District Collectorates, Municipal Corporations & Smart Cities.\n• Private Sector & Public: Environmental consultancies, healthcare providers, schools, logistics fleets, and citizens.",
         ACCENT_BLUE),
        ("💰 Business & Deployment Model",
         "• B2G SaaS & Enterprise API: Tiered subscription model for municipal administrations & state boards.\n• Zero-Capex Deployment: Low operational cost on cloud infrastructure; 99% cheaper than physical ground sensors.\n• Public Good API: Open freemium citizen tier for community air health alerts.",
         ACCENT_GREEN),
        ("📈 Market Size & Scalability",
         "• Global Environmental Monitoring Market: $3.2 Billion (growing at 8.4% CAGR).\n• Pan-India Scalability: Modular coordinate grid easily expands from 15 pilot districts to all 700+ districts across India.\n• Global Portability: The satellite-first architecture can be deployed in Southeast Asia, Latin America, and Sub-Saharan Africa with zero code changes.",
         ACCENT_AMBER)
    ]

    for i, (mtitle, mbody, mcolor) in enumerate(market_points):
        left = Inches(0.8 + i * 3.95)
        card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.6), Inches(3.75), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = mcolor
        card.line.width = Pt(1.5)

        mtf = card.text_frame
        mtf.word_wrap = True
        p1 = mtf.paragraphs[0]
        p1.text = mtitle
        p1.font.size = Pt(15)
        p1.font.bold = True
        p1.font.color.rgb = mcolor

        p2 = mtf.add_paragraph()
        p2.text = "\n" + mbody
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = TEXT_MUTED

    # Save presentation
    output_path = "VayuShetra_SIH_Official_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation generated successfully at {output_path}")

if __name__ == "__main__":
    create_presentation()
